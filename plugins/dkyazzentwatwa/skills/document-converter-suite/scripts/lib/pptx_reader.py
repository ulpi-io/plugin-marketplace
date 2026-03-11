from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional

from pptx import Presentation


@dataclass
class SlideContent:
    index: int
    title: str
    texts: List[str]
    tables: List[List[List[str]]]


@dataclass
class PptxContent:
    slides: List[SlideContent]


def _shape_text(shape) -> Optional[str]:
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            txt = shape.text_frame.text
            return (txt or "").strip()
    except Exception:
        return None
    return None


def read_pptx_content(path: Path, max_chars: int = 300000) -> PptxContent:
    """Extract slide titles, text shapes, and tables from a .pptx."""
    prs = Presentation(str(path))
    slides: List[SlideContent] = []

    budget = max_chars
    for i, slide in enumerate(prs.slides):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except Exception:
            title = ""

        texts: List[str] = []
        tables: List[List[List[str]]] = []

        for shape in slide.shapes:
            # tables
            try:
                if hasattr(shape, "has_table") and shape.has_table:
                    t = shape.table
                    rows: List[List[str]] = []
                    for r in t.rows:
                        rows.append([(c.text or "").strip() for c in r.cells])
                    if rows:
                        tables.append(rows)
                    continue
            except Exception:
                pass

            txt = _shape_text(shape)
            if not txt:
                continue
            if txt == title:
                continue
            if budget <= 0:
                break
            if len(txt) > budget:
                txt = txt[:budget]
            budget -= len(txt)
            texts.append(txt)

        slides.append(SlideContent(index=i + 1, title=title, texts=texts, tables=tables))

    return PptxContent(slides=slides)
