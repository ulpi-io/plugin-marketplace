from __future__ import annotations

from pathlib import Path
from typing import List, Optional

from pptx import Presentation


def write_pptx_from_slides(
    output_path: Path,
    slides: List[dict],
    title: Optional[str] = None,
) -> None:
    """Write a basic PPTX.

    Each slide dict can have:
      - title: str
      - bullets: list[str]
      - tables: list[list[list[str]]]
    """
    prs = Presentation()

    # Optional title slide
    if title:
        try:
            layout = prs.slide_layouts[0]
        except Exception:
            layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title

    for s in slides:
        try:
            layout = prs.slide_layouts[1]  # Title + Content
        except Exception:
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = (s.get("title") or "").strip()

        # Text / bullets into the body placeholder if present
        body = None
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type == 2:  # BODY
                    body = shape
                    break
            except Exception:
                continue
        if body is None and len(slide.shapes) > 1:
            # best-effort: use second shape if it's a placeholder
            try:
                body = slide.shapes[1]
            except Exception:
                body = None

        bullets = s.get("bullets") or []
        if body is not None and hasattr(body, "text_frame"):
            tf = body.text_frame
            tf.clear()
            if bullets:
                first = tf.paragraphs[0]
                first.text = bullets[0]
                first.level = 0
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
            else:
                tf.text = ""

        # Tables: create one slide per table when multiple tables exist
        tables = s.get("tables") or []
        if tables:
            # Render first table on the current slide
            for table_idx, t in enumerate(tables):
                # For tables after the first, create new slides
                if table_idx > 0:
                    try:
                        layout = prs.slide_layouts[1]  # Title + Content
                    except Exception:
                        layout = prs.slide_layouts[0]

                    slide = prs.slides.add_slide(layout)
                    slide_title = s.get("title") or ""
                    if slide.shapes.title:
                        slide.shapes.title.text = f"{slide_title} - Table {table_idx + 1}"

                if not t:
                    continue

                rows = len(t)
                cols = max((len(r) for r in t), default=0)
                if rows <= 0 or cols <= 0:
                    continue

                # Simple positioning: below title, centered-ish
                try:
                    from pptx.util import Inches

                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(5)
                    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                    tbl = shape.table
                except Exception:
                    continue

                for r_i in range(rows):
                    for c_i in range(cols):
                        val = t[r_i][c_i] if c_i < len(t[r_i]) else ""
                        tbl.cell(r_i, c_i).text = str(val or "")

    prs.save(str(output_path))
